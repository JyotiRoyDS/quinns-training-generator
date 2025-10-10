import streamlit as st
import os
import io
import re
import json
import time
import base64
import tempfile
import shutil
from datetime import datetime, timedelta
from pathlib import Path
from typing import List, Dict, Tuple, Any, Optional, Union
import asyncio
import logging
from concurrent.futures import ThreadPoolExecutor

# Document processing
import pdfplumber
import pytesseract
from pdf2image import convert_from_path
import docx
from pptx import Presentation
import pandas as pd

# Basic NLP
from sklearn.feature_extraction.text import CountVectorizer
from sklearn.decomposition import LatentDirichletAllocation
import numpy as np

# Document generation
from docx import Document as DocxDocument
from docx.shared import RGBColor as DocxRGBColor, Pt as DocxPt, Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# API - OpenAI 1.0+ compatible
from openai import OpenAI
from tenacity import retry, wait_exponential, stop_after_attempt, retry_if_exception_type
def get_api_key():
    """Load API key from Streamlit secrets or environment variable"""
    # Try Streamlit secrets first (for cloud deployment)
    try:
        if hasattr(st, 'secrets') and 'OPENAI_API_KEY' in st.secrets:
            api_key = st.secrets['OPENAI_API_KEY']
            if api_key and api_key.strip():
                return api_key.strip()
    except Exception:
        # Logger not available yet at this point
        pass
    
    # Fall back to environment variable (for local development)
    api_key = os.getenv('OPENAI_API_KEY')
    if api_key and api_key.strip():
        return api_key.strip()
    
    # Return None if no key found
    return None

OPENAI_API_KEY = get_api_key()
# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger("QUINNS-TrainingGenerator")

# Constants
MAX_WORKERS = 3
DEFAULT_TOPICS = 5
DEFAULT_SLIDES_PER_TOPIC = 5

# Training duration constants (in minutes per slide)
SLIDE_DURATION_MAP = {
    "fast": 3,      # 3 minutes per slide
    "medium": 5,    # 5 minutes per slide  
    "thorough": 8   # 8 minutes per slide
}

# Create output directories
try:
    TEMP_DIR = Path.cwd() / "quinns_output"
    os.makedirs(TEMP_DIR, exist_ok=True)
    
    CACHE_DIR = TEMP_DIR / "cache"
    os.makedirs(CACHE_DIR, exist_ok=True)
    
    OUTPUT_DIR = TEMP_DIR / "output"
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    TEMPLATES_DIR = TEMP_DIR / "templates"
    os.makedirs(TEMPLATES_DIR, exist_ok=True)
    
    logger.info(f"Created output directories at {TEMP_DIR}")
except Exception as e:
    logger.error(f"Failed to create output directories: {e}")
    TEMP_DIR = Path(tempfile.gettempdir()) / "quinns_generator"
    CACHE_DIR = TEMP_DIR / "cache"
    OUTPUT_DIR = TEMP_DIR / "output"
    TEMPLATES_DIR = TEMP_DIR / "templates"
    os.makedirs(TEMP_DIR, exist_ok=True)
    os.makedirs(CACHE_DIR, exist_ok=True)
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(TEMPLATES_DIR, exist_ok=True)

# ENHANCED PROMPT TEMPLATES
PROMPTS = {
    "topic_extraction": """
    Analyze this training document and extract the main topics covered.
    Return a JSON object with specific, concrete topics found in the document.
    
    DOCUMENT CONTENT:
    {content}
    
    Return format:
    {{
        "topics": [
            {{
                "title": "Specific topic name from document",
                "description": "What this topic actually covers in the document",
                "importance": "high/medium/low",
                "key_concepts": ["concept1", "concept2", "concept3"],
                "estimated_duration_minutes": 30
            }}
        ]
    }}
    """,
    
    "outline_generation": """
    Create a detailed training outline based on the following content.
    Extract real information from the document to create specific, actionable learning objectives and content.
    
    TOPICS IDENTIFIED:
    {topics}
    
    DOCUMENT CONTENT:
    {content}
    
    TRAINING DURATION: {duration}
    
    Return a JSON object with this structure:
    {{
        "title": "Descriptive training title based on content",
        "description": "What this training actually teaches",
        "target_audience": "Who needs this training",
        "prerequisites": "What learners should know beforehand",
        "difficulty_level": "beginner/intermediate/advanced",
        "estimated_duration_hours": 8,
        "objectives": [
            "Specific, measurable learning outcome 1",
            "Specific, measurable learning outcome 2"
        ],
        "modules": [
            {{
                "title": "Module title from content",
                "duration": "30-60 minutes",
                "difficulty": "beginner/intermediate/advanced",
                "objectives": ["What learners will achieve in this module"],
                "topics": ["Specific topics covered"],
                "key_points": ["Actual facts/concepts from document"],
                "activities": ["Specific practice activities"],
                "estimated_slides": 8
            }}
        ]
    }}
    
    Focus on extracting REAL information from the document, not generic placeholders.
    """,
    
    "slide_generation": """
    Generate {num_slides} slides with specific content from the source material.

    TOPIC: {topic}
    NUMBER OF SLIDES: {num_slides}
    DIFFICULTY: {difficulty}

    SOURCE CONTEXT:
    {context}

    ALREADY COVERED (avoid repeating):
    {previous_context}

    INSTRUCTIONS:
    1. Extract real facts, procedures, and examples from SOURCE CONTEXT
    2. Each slide covers DIFFERENT information
    3. Write complete sentences with specific details
    4. Include examples and real-world applications
    5. Notes must be 150+ words with teaching content

    Return a JSON object with this exact structure:
    {{
        "slides": [
            {{
                "slide_type": "content",
                "title": "Specific topic title",
                "content": [
                    "Point 1 with specific information",
                    "Point 2 with details",
                    "Point 3 with examples",
                    "Point 4 with key takeaway"
                ],
                "notes": "Detailed teaching notes with specific information, examples, and explanations from the source material. Minimum 150 words.",
                "estimated_time_minutes": 5
            }}
        ]
    }}
    
    IMPORTANT: Return ONLY valid JSON with a "slides" array. No extra text.
    """,
    
    "assessment_generation": """
    Create comprehensive assessment questions based on the training content.
    
    MODULE CONTENT:
    {module_content}
    
    LEARNING OBJECTIVES:
    {objectives}
    
    DIFFICULTY LEVEL: {difficulty}
    NUMBER OF QUESTIONS: {num_questions}
    
    Generate a mix of question types that assess understanding of the material.
    
    Return JSON:
    {{
        "questions": [
            {{
                "type": "multiple_choice",
                "question": "Specific question based on content",
                "options": ["Option A", "Option B", "Option C", "Option D"],
                "correct_answer": "Option B",
                "explanation": "Why this is correct and others are wrong",
                "difficulty": "easy/medium/hard",
                "learning_objective": "Which objective this assesses"
            }},
            {{
                "type": "true_false",
                "question": "Statement to evaluate",
                "correct_answer": true,
                "explanation": "Explanation of the correct answer"
            }},
            {{
                "type": "scenario",
                "scenario": "Real-world situation description",
                "question": "What should you do?",
                "suggested_answer": "Detailed response with key points",
                "rubric": ["Criterion 1", "Criterion 2", "Criterion 3"]
            }}
        ]
    }}
    """,
    
    "activity_generation": """
    Create engaging learning activities based on the module content.
    
    MODULE: {module_title}
    CONTENT: {module_content}
    DURATION: {duration} minutes
    
    Generate practical, hands-on activities that reinforce learning.
    
    Return JSON:
    {{
        "activities": [
            {{
                "type": "case_study",
                "title": "Activity title",
                "description": "What participants will do",
                "scenario": "Detailed realistic scenario",
                "instructions": ["Step 1", "Step 2", "Step 3"],
                "discussion_questions": ["Question 1", "Question 2"],
                "duration_minutes": 15,
                "group_size": "2-4 people"
            }},
            {{
                "type": "role_play",
                "title": "Activity title",
                "roles": ["Role 1", "Role 2"],
                "scenario": "Situation to act out",
                "objectives": ["What to demonstrate"],
                "debrief_questions": ["Reflection question"]
            }}
        ]
    }}
    """,
    
    "quality_analysis": """
    Analyze the quality and completeness of this training content.
    
    CONTENT:
    {content}
    
    OUTLINE:
    {outline}
    
    Evaluate on:
    - Content clarity and completeness
    - Learning objective alignment
    - Appropriate difficulty level
    - Engagement potential
    - Practical applicability
    
    Return JSON:
    {{
        "overall_score": 85,
        "readability_score": 75,
        "completeness_score": 90,
        "engagement_score": 80,
        "strengths": ["Strength 1", "Strength 2"],
        "improvements": ["Suggestion 1", "Suggestion 2"],
        "missing_topics": ["Topic that should be covered"],
        "recommendations": ["Recommendation 1", "Recommendation 2"]
    }}
    """,
    
    "trainer_guide": """
    Create COMPREHENSIVE trainer guide content with ACTUAL teaching material.

    MODULE CONTEXT: {module}
    SLIDE: {slide}

    Generate a detailed trainer guide section with:

    1. WHAT TO TEACH: The actual content (facts, procedures, concepts)
    2. HOW TO TEACH IT: Delivery methods and techniques
    3. EXAMPLES: Real examples from the source material
    4. PRACTICE: Hands-on exercises or discussions
    5. Q&A: Likely questions with complete answers

    Return JSON:
    {{
        "title": "Section title from slide",
        "teaching_content": "3-4 paragraphs of the ACTUAL information the trainer should teach - include all facts, procedures, definitions, examples from the source material. This should be substantial enough that a trainer could teach from this alone. Minimum 300 words with specific details, not generic guidance.",
        "delivery_approach": "Specific teaching methods for this content (lecture, demonstration, discussion, etc.) with timing breakdown",
        "real_examples": [
            "Specific example 1 with full context and details from source material",
            "Specific example 2 from source material with concrete information",
            "Case study or scenario with complete information and real data"
        ],
        "key_points_to_emphasize": [
            "Critical fact/concept with detailed explanation",
            "Important procedure/process with complete step-by-step breakdown",
            "Common misconception to address with correction"
        ],
        "interactive_activity": {{
            "activity_type": "Type of activity (discussion, exercise, demo, etc.)",
            "instructions": "Complete step-by-step instructions with specific details",
            "duration": "X minutes",
            "materials": ["specific", "materials", "needed"]
        }},
        "anticipated_questions": [
            {{
                "question": "Specific question learners might ask",
                "detailed_answer": "Complete, thorough answer with examples and explanation (minimum 100 words)"
            }},
            {{
                "question": "Another likely question",
                "detailed_answer": "Another detailed response"
            }}
        ],
        "assessment_checkpoint": "Quick knowledge check question or exercise with answer",
        "timing_guidance": "X-Y minutes with breakdown of time allocation",
        "transitions": "How to transition from previous slide and set up next slide"
    }}

    CRITICAL: This must contain ACTUAL teaching content with real information, not meta-instructions about what to teach.
    """
}

# TRAINING TEMPLATES
TRAINING_TEMPLATES = {
    "corporate": {
        "name": "Corporate Training",
        "description": "Professional business training format",
        "colors": {
            "primary": "4F46E5",
            "secondary": "10B981",
            "accent": "F59E0B"
        },
        "style": "formal",
        "includes_assessments": True,
        "includes_activities": True
    },
    "technical": {
        "name": "Technical Skills",
        "description": "Hands-on technical training",
        "colors": {
            "primary": "3B82F6",
            "secondary": "8B5CF6",
            "accent": "EC4899"
        },
        "style": "detailed",
        "includes_assessments": True,
        "includes_activities": True
    },
    "compliance": {
        "name": "Compliance & Safety",
        "description": "Regulatory and safety training",
        "colors": {
            "primary": "EF4444",
            "secondary": "F59E0B",
            "accent": "10B981"
        },
        "style": "structured",
        "includes_assessments": True,
        "includes_activities": False
    },
    "onboarding": {
        "name": "New Employee Onboarding",
        "description": "Welcome and orientation training",
        "colors": {
            "primary": "10B981",
            "secondary": "3B82F6",
            "accent": "F59E0B"
        },
        "style": "welcoming",
        "includes_assessments": False,
        "includes_activities": True
    },
    "leadership": {
        "name": "Leadership Development",
        "description": "Management and leadership skills",
        "colors": {
            "primary": "7C3AED",
            "secondary": "EC4899",
            "accent": "F59E0B"
        },
        "style": "inspirational",
        "includes_assessments": True,
        "includes_activities": True
    }
}

def get_cache_key(file_path: Union[str, Path], process_type: str) -> str:
    if isinstance(file_path, Path):
        file_path = str(file_path)
    if os.path.isfile(file_path):
        file_hash = str(hash(file_path + str(os.path.getmtime(file_path))))
    else:
        file_hash = str(hash(file_path))
    return f"{file_hash}_{process_type}"

def get_from_cache(cache_key: str) -> Optional[Dict]:
    cache_file = CACHE_DIR / f"{cache_key}.json"
    if cache_file.exists():
        try:
            with open(cache_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"Error reading cache: {e}")
    return None

def save_to_cache(cache_key: str, data: Dict) -> None:
    cache_file = CACHE_DIR / f"{cache_key}.json"
    try:
        with open(cache_file, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"Error saving to cache: {e}")

def split_into_sentences(text: str) -> List[str]:
    sentences = []
    for paragraph in text.split('\n'):
        if paragraph.strip():
            if re.match(r'^[\s]*[•\-\*\d]+[\.\)]*\s', paragraph):
                sentences.append(paragraph.strip())
            else:
                for s in re.split(r'(?<=[.!?])\s+', paragraph):
                    if len(s.strip().split()) > 3:
                        sentences.append(s.strip())
    return sentences

def calculate_training_duration(num_slides: int, pace: str = "medium", includes_activities: bool = True) -> Dict[str, float]:
    """Calculate training duration based on slides and pace"""
    minutes_per_slide = SLIDE_DURATION_MAP.get(pace, 5)
    
    # Base time for slides
    presentation_time = num_slides * minutes_per_slide
    
    # Add time for activities (roughly 20% more time)
    if includes_activities:
        presentation_time *= 1.2
    
    # Add breaks (10 min per hour)
    total_hours = presentation_time / 60
    break_time = (total_hours // 1) * 10  # 10 min break per hour
    
    total_minutes = presentation_time + break_time
    
    return {
        "total_minutes": total_minutes,
        "total_hours": total_minutes / 60,
        "days": total_minutes / (60 * 8),  # 8-hour days
        "presentation_minutes": presentation_time,
        "break_minutes": break_time
    }

def format_duration_display(duration_info: Dict[str, float]) -> str:
    """Format duration for display"""
    hours = int(duration_info['total_hours'])
    minutes = int(duration_info['total_minutes'] % 60)
    days = duration_info['days']
    
    if days >= 1:
        return f"{days:.1f} days ({hours}h {minutes}m)"
    elif hours >= 1:
        return f"{hours}h {minutes}m"
    else:
        return f"{minutes}m"

class DocumentProcessor:
    """Handles extraction of text from various document formats"""
    
    @staticmethod
    def process_document(file_path: Union[str, Path]) -> str:
        if isinstance(file_path, str):
            file_path = Path(file_path)
            
        cache_key = get_cache_key(file_path, "text_extraction")
        cached_content = get_from_cache(cache_key)
        
        if cached_content:
            logger.info(f"Using cached extraction for {file_path.name}")
            return cached_content["content"]
        
        logger.info(f"Processing document: {file_path.name}")
        
        ext = file_path.suffix.lower()
        content = ""
        
        try:
            if ext == ".pdf":
                content = DocumentProcessor._process_pdf(file_path)
            elif ext == ".docx":
                content = DocumentProcessor._process_docx(file_path)
            elif ext == ".pptx":
                content = DocumentProcessor._process_pptx(file_path)
            elif ext in [".txt", ".md"]:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    content = f.read()
            else:
                raise ValueError(f"Unsupported file format: {ext}")
                
            content = re.sub(r'\s+', ' ', content).strip()
            save_to_cache(cache_key, {"content": content, "timestamp": datetime.now().isoformat()})
            
            return content
        except Exception as e:
            logger.error(f"Error processing document {file_path.name}: {e}")
            raise
    
    @staticmethod
    def _process_pdf(file_path: Path) -> str:
        try:
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text() or ""
                    text += page_text + "\n\n"
            
            if len(text.strip()) < 100:
                logger.info(f"PDF appears to be scanned, applying OCR: {file_path.name}")
                try:
                    images = convert_from_path(file_path)
                    ocr_text = ""
                    for img in images:
                        ocr_text += pytesseract.image_to_string(img) + "\n\n"
                    if len(ocr_text.strip()) > len(text.strip()):
                        text = ocr_text
                except Exception as ocr_error:
                    logger.error(f"OCR processing failed: {ocr_error}")
            
            return text
        except Exception as e:
            logger.error(f"PDF processing error: {e}")
            raise
    
    @staticmethod
    def _process_docx(file_path: Path) -> str:
        try:
            doc = docx.Document(file_path)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text)
                    text += " | ".join(row_text) + "\n"
                text += "\n"
            return text
        except Exception as e:
            logger.error(f"DOCX processing error: {e}")
            raise
    
    @staticmethod
    def _process_pptx(file_path: Path) -> str:
        try:
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                title = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        title = shape.text
                        break
                text += f"Slide {slide_num}: {title}\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text != title:
                        text += shape.text + "\n"
                try:
                    if hasattr(slide, "notes_slide") and slide.notes_slide:
                        if hasattr(slide.notes_slide, "notes_text_frame") and slide.notes_slide.notes_text_frame.text:
                            text += "Notes: " + slide.notes_slide.notes_text_frame.text + "\n"
                except Exception as note_error:
                    logger.warning(f"Error extracting notes from slide {slide_num}: {note_error}")
                text += "\n"
            return text
        except Exception as e:
            logger.error(f"PPTX processing error: {e}")
            raise

class ContentAnalyzer:
    """Analyzes document content with quality metrics"""
    
    def __init__(self, openai_api_key: str):
        self.client = OpenAI(api_key=openai_api_key)
        logger.info("ContentAnalyzer initialized with OpenAI API")
    
    def extract_topics(self, content: str) -> List[Dict]:
        cache_key = get_cache_key(str(hash(content)), "topic_extraction")
        cached_topics = get_from_cache(cache_key)
        
        if cached_topics:
            logger.info("Using cached topic extraction")
            return cached_topics["topics"]
        
        try:
            topics = self._extract_topics_openai(content)
            save_to_cache(cache_key, {"topics": topics, "timestamp": datetime.now().isoformat()})
            return topics
        except Exception as e:
            logger.error(f"OpenAI topic extraction error: {e}")
            return self._fallback_topic_extraction(content)
    
    @retry(wait=wait_exponential(multiplier=1, min=4, max=60), stop=stop_after_attempt(3))
    def _extract_topics_openai(self, content: str) -> List[Dict]:
        content_sample = self._smart_sample(content, 6000)  # Increased from 4000
        
        prompt = PROMPTS["topic_extraction"].format(content=content_sample)
        
        response = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are an expert at analyzing training documents and extracting specific, concrete topics with real information from the content."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.2,
            response_format={"type": "json_object"},
            max_tokens=1500
        )
        
        result = response.choices[0].message.content
        topics_data = json.loads(result)
        
        if "topics" not in topics_data:
            raise ValueError("Invalid response format from OpenAI")
            
        return topics_data["topics"]
    
    def _smart_sample(self, text: str, target_chars: int) -> str:
        """Sample text more intelligently to preserve context - ENHANCED"""
        # DOUBLED the context window
        target_chars = target_chars * 2
        
        if len(text) <= target_chars:
            return text
            
        sentences = split_into_sentences(text)
        if not sentences:
            return text[:target_chars]
        
        # Take MORE from beginning and end
        first_count = int(len(sentences) * 0.5)  # 50% from start (was 40%)
        last_count = int(len(sentences) * 0.3)   # 30% from end (was 20%)
        middle_count = len(sentences) - first_count - last_count
        
        selected = sentences[:first_count]
        
        if middle_count > 0:
            stride = max(1, middle_count // 30)  # Sample more from middle (was 20)
            for i in range(first_count, first_count + middle_count, stride):
                if i < len(sentences):
                    selected.append(sentences[i])
        
        selected.extend(sentences[-last_count:])
        
        result = " ".join(selected)
        return result[:target_chars * 2] if len(result) > target_chars * 2 else result
    
    def _fallback_topic_extraction(self, content: str) -> List[Dict]:
        sentences = split_into_sentences(content)
        
        if len(sentences) < 10:
            words = content.lower().split()
            common_words = set(['the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'from', 'as', 'is', 'was', 'are', 'were'])
            key_words = [w for w in words if w not in common_words and len(w) > 4][:20]
            
            return [{
                "title": " ".join(key_words[:3]).title(),
                "description": sentences[0] if sentences else "Main content topic",
                "importance": "high",
                "key_concepts": key_words[:5],
                "estimated_duration_minutes": 30
            }]
        
        try:
            vectorizer = CountVectorizer(max_df=0.9, min_df=2, stop_words='english', max_features=5000)
            word_counts = vectorizer.fit_transform(sentences)
            
            lda = LatentDirichletAllocation(n_components=min(5, len(sentences)//10), max_iter=10, random_state=42)
            lda.fit(word_counts)
            
            feature_names = vectorizer.get_feature_names_out()
            topics = []
            
            for i, topic in enumerate(lda.components_):
                top_words_idx = topic.argsort()[:-8-1:-1]
                top_words = [feature_names[idx] for idx in top_words_idx]
                
                topics.append({
                    "title": " ".join(top_words[:2]).title(),
                    "description": f"Covers: {', '.join(top_words[:5])}",
                    "importance": "medium",
                    "key_concepts": top_words[:5],
                    "estimated_duration_minutes": 30
                })
                
            return topics
        except Exception as e:
            logger.error(f"Fallback topic extraction failed: {e}")
            return [{"title": "Document Content", "description": "Main topics from document", "importance": "high", "key_concepts": [], "estimated_duration_minutes": 30}]
    
    @retry(wait=wait_exponential(multiplier=1, min=4, max=60), stop=stop_after_attempt(3))
    def detect_training_type(self, content: str) -> str:
        """Auto-detect the type of training from content"""
        content_sample = self._smart_sample(content, 2000)
        
        prompt = f"""
        Analyze this training content and determine the most appropriate training type.
        
        CONTENT:
        {content_sample}
        
        Choose from these types:
        - Corporate Training: General business skills, professional development
        - Technical Skills: IT, software, technical procedures
        - Compliance & Safety: Regulations, safety protocols, legal requirements
        - New Employee Onboarding: Company introduction, orientation
        - Leadership Development: Management, leadership skills
        - Sales Training: Sales techniques, customer acquisition
        - Customer Service: Service excellence, customer interaction
        - Soft Skills: Communication, teamwork, emotional intelligence
        - Product Training: Specific product knowledge
        
        Return JSON:
        {{
            "training_type": "The most appropriate type",
            "confidence": 0.85,
            "reasoning": "Why this type was selected"
        }}
        """
        
        response = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a training classification expert."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.2,
            response_format={"type": "json_object"},
            max_tokens=300
        )
        
        result = response.choices[0].message.content
        detection_result = json.loads(result)
        return detection_result.get("training_type", "Corporate Training")
    
    @retry(wait=wait_exponential(multiplier=1, min=4, max=60), stop=stop_after_attempt(3))
    def generate_outline(self, topics: List[Dict], content: str, duration: str = "1 day") -> Dict:
        cache_key = get_cache_key(str(hash(str(topics) + content)), f"outline_{duration}")
        cached_outline = get_from_cache(cache_key)
        
        if cached_outline:
            logger.info("Using cached outline")
            return cached_outline["outline"]
        
        topic_summary = "\n".join([
            f"- {topic['title']}: {topic['description']}\n  Key concepts: {', '.join(topic.get('key_concepts', []))}"
            for topic in topics
        ])
        
        content_summary = self._smart_sample(content, 6000)  # Increased from 4000
        
        prompt = PROMPTS["outline_generation"].format(
            topics=topic_summary,
            content=content_summary,
            duration=duration
        )
        
        response = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a professional training developer. Extract REAL, SPECIFIC information from the provided document to create detailed learning content. Do not write generic placeholders."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            response_format={"type": "json_object"},
            max_tokens=2500
        )
        
        result = response.choices[0].message.content
        outline = json.loads(result)
        save_to_cache(cache_key, {"outline": outline, "timestamp": datetime.now().isoformat()})
        
        return outline
    
    @retry(wait=wait_exponential(multiplier=1, min=4, max=60), stop=stop_after_attempt(3))
    def analyze_quality(self, content: str, outline: Dict) -> Dict:
        """Analyze content quality"""
        prompt = PROMPTS["quality_analysis"].format(
            content=self._smart_sample(content, 3000),
            outline=json.dumps(outline, indent=2)[:2000]
        )
        
        response = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a training quality analyst. Provide honest, constructive feedback."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            response_format={"type": "json_object"},
            max_tokens=1000
        )
        
        result = response.choices[0].message.content
        return json.loads(result)

class AssessmentGenerator:
    """Generates assessments and quizzes"""
    
    def __init__(self, openai_api_key: str):
        self.client = OpenAI(api_key=openai_api_key)
    
    @retry(wait=wait_exponential(multiplier=1, min=4, max=60), stop=stop_after_attempt(3))
    async def generate_assessment(self, module: Dict, num_questions: int = 10) -> Dict:
        """Generate assessment questions for a module"""
        module_content = json.dumps({
            "title": module.get("title"),
            "key_points": module.get("key_points", []),
            "topics": module.get("topics", [])
        }, indent=2)
        
        objectives = "\n".join(module.get("objectives", []))
        difficulty = module.get("difficulty", "intermediate")
        
        prompt = PROMPTS["assessment_generation"].format(
            module_content=module_content,
            objectives=objectives,
            difficulty=difficulty,
            num_questions=num_questions
        )
        
        response = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are an assessment design expert. Create fair, clear, and relevant questions that accurately assess learning."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.4,
            response_format={"type": "json_object"},
            max_tokens=2000
        )
        
        result = response.choices[0].message.content
        return json.loads(result)
    
    @retry(wait=wait_exponential(multiplier=1, min=4, max=60), stop=stop_after_attempt(3))
    async def generate_activities(self, module: Dict) -> Dict:
        """Generate learning activities"""
        duration = module.get("duration", "45 minutes")
        duration_minutes = int(re.search(r'\d+', duration).group()) if re.search(r'\d+', duration) else 45
        
        prompt = PROMPTS["activity_generation"].format(
            module_title=module.get("title"),
            module_content=json.dumps(module, indent=2)[:1500],
            duration=duration_minutes
        )
        
        response = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a learning experience designer. Create engaging, practical activities that reinforce learning."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.5,
            response_format={"type": "json_object"},
            max_tokens=1500
        )
        
        result = response.choices[0].message.content
        return json.loads(result)

class DocumentGenerator:
    """Generates training materials with brand customization"""
    
    def __init__(self, openai_api_key: str, brand_config: Optional[Dict] = None):
        self.client = OpenAI(api_key=openai_api_key)
        self.content_store = {}
        self.brand_config = brand_config or {}
        self.assessment_generator = AssessmentGenerator(openai_api_key)
    
    def set_source_content(self, content: str):
        self.content_store["source"] = content
        self.content_store["sentences"] = split_into_sentences(content)
    
    def set_brand_config(self, brand_config: Dict):
        self.brand_config = brand_config
    
    def _get_relevant_context(self, topic: str, num_sentences: int = 30) -> str:
        """Get MORE relevant context with better scoring - ENHANCED"""
        if "sentences" not in self.content_store:
            return ""
        
        sentences = self.content_store["sentences"]
        topic_lower = topic.lower()
        topic_words = set(topic_lower.split())
        
        # Score sentences by relevance
        scored = []
        for sent in sentences:
            sent_lower = sent.lower()
            
            # Multiple scoring factors
            word_matches = sum(1 for word in topic_words if word in sent_lower)
            length_bonus = 1 if len(sent.split()) > 15 else 0  # Prefer detailed sentences
            has_numbers = 1 if any(char.isdigit() for char in sent) else 0  # Facts with numbers
            has_proper_nouns = 1 if any(word[0].isupper() for word in sent.split()[1:]) else 0
            
            score = (word_matches * 2) + length_bonus + has_numbers + has_proper_nouns
            
            if score > 0:
                scored.append((score, sent))
        
        scored.sort(reverse=True)
        
        # Get top sentences PLUS surrounding context for continuity
        relevant = []
        added_indices = set()
        
        for score, sent in scored[:num_sentences]:
            idx = sentences.index(sent)
            
            # Add the sentence itself
            if idx not in added_indices:
                relevant.append(sent)
                added_indices.add(idx)
            
            # Add next sentence for context continuity
            if idx + 1 < len(sentences) and (idx + 1) not in added_indices:
                relevant.append(sentences[idx + 1])
                added_indices.add(idx + 1)
        
        return " ".join(relevant[:num_sentences * 2])  # Allow more context
    
    def _validate_slide_content(self, slide: Dict) -> bool:
        """Check if slide has substantive, unique content"""
        content = slide.get('content', [])
        notes = slide.get('notes', '')
        
        # Check for generic placeholders
        generic_phrases = [
            'key concept', 'important topic', 'discuss the', 'present information',
            'cover the basics', 'introduce the', 'explain the importance', 'overview of'
        ]
        
        content_text = " ".join(str(c) for c in content).lower()
        
        # Reject if too generic
        if any(phrase in content_text for phrase in generic_phrases):
            logger.warning(f"Rejected generic slide: {slide.get('title')}")
            return False
        
        # Reject if notes too short
        if len(notes) < 100:
            logger.warning(f"Rejected slide with short notes: {slide.get('title')}")
            return False
        
        # Require some specificity (numbers, proper nouns, technical terms)
        has_specificity = (
            any(char.isdigit() for char in content_text) or
            any(word[0].isupper() for word in content_text.split()[1:]) or
            len(content_text) > 200
        )
        
        if not has_specificity:
            logger.warning(f"Rejected non-specific slide: {slide.get('title')}")
            return False
        
        return True
    
    @retry(wait=wait_exponential(multiplier=1, min=4, max=60), stop=stop_after_attempt(3), 
           retry=retry_if_exception_type((Exception,)))
    async def generate_slides_batch(self, topic: str, module: Dict, num_slides: int, previous_context: str = "") -> List[Dict]:
        try:
            # Get MORE context - increased significantly
            context = self._get_relevant_context(topic, num_sentences=50)  # Was 10
            
            if not context:
                context = " ".join(module.get("key_points", [topic]))
            
            # Limit context size to avoid token issues
            max_context_chars = 4000
            if len(context) > max_context_chars:
                context = context[:max_context_chars]
            
            if len(previous_context) > 2000:
                previous_context = previous_context[:2000]
            
            difficulty = module.get("difficulty", "intermediate")
            
            prompt = PROMPTS["slide_generation"].format(
                topic=topic,
                num_slides=num_slides,
                difficulty=difficulty,
                context=context,
                previous_context=previous_context
            )
            
            # Log prompt size for debugging
            logger.info(f"Generating {num_slides} slides for topic: {topic}")
            logger.info(f"Prompt length: {len(prompt)} chars")
            
            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "You are creating training slide content. Write ACTUAL educational content with specific facts, procedures, and examples from the SOURCE CONTEXT - NOT generic instructions about what to present. Be concrete and specific. Extract real information. Return a JSON object with a 'slides' array."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,
                response_format={"type": "json_object"},
                max_tokens=min(4000, 1000 * num_slides)  # Cap at 4000 tokens
            )
            
            result = response.choices[0].message.content
            logger.info(f"Received response for {topic}")
            
            try:
                slides_data = json.loads(result)
            except json.JSONDecodeError as e:
                logger.error(f"JSON decode error: {e}")
                logger.error(f"Response content: {result[:500]}")
                raise
            
            # Handle both list and dict formats
            if isinstance(slides_data, list):
                slides_list = slides_data
            elif isinstance(slides_data, dict):
                slides_list = slides_data.get("slides", [])
                if not slides_list and "slide" in slides_data:
                    slides_list = [slides_data]
            else:
                logger.error(f"Unexpected slides_data type: {type(slides_data)}")
                slides_list = []
            
            # Validate slides before returning
            validated_slides = []
            for slide in slides_list:
                if isinstance(slide, dict) and self._validate_slide_content(slide):
                    validated_slides.append(slide)
            
            if not validated_slides:
                logger.warning(f"No validated slides for {topic}, returning original")
                return slides_list if slides_list else []
            
            return validated_slides
            
        except Exception as e:
            logger.error(f"Error in generate_slides_batch for {topic}: {e}")
            logger.error(f"Error type: {type(e).__name__}")
            # Return empty list instead of raising
            return []
    
    def _create_fallback_slide(self, topic: str, module: Dict = None) -> Dict:
        content = []
        notes = ""
        
        if module and "key_points" in module:
            content = module["key_points"][:4]
            notes = f"{topic}: " + " ".join(module["key_points"])
        else:
            context = self._get_relevant_context(topic, 10)
            if context:
                sentences = split_into_sentences(context)
                content = sentences[:4]
                notes = context
            else:
                content = [
                    f"{topic} - Key Concept #1",
                    f"{topic} - Key Concept #2",
                    f"{topic} - Key Concept #3",
                    f"{topic} - Key Concept #4"
                ]
                notes = f"Core information about {topic}"
        
        return {
            "slide_type": "content",
            "title": topic,
            "content": content,
            "notes": notes,
            "estimated_time_minutes": 5
        }
    
    async def generate_slides_for_module(self, module: Dict, max_slides: int = 8) -> List[Dict]:
        """Generate slides for a module with better context tracking and rate limiting - ENHANCED"""
        slides = []
        
        # Title slide with actual content
        title_slide = {
            "slide_type": "title",
            "title": module["title"],
            "content": module["objectives"],
            "notes": f"Welcome to {module['title']}. " + " ".join(module.get('key_points', [])[:3]),
            "estimated_time_minutes": 2
        }
        slides.append(title_slide)
        
        # Track covered content to avoid repetition
        covered_topics = set()
        context_history = f"Module: {module['title']}\nObjectives: {', '.join(module['objectives'])}\n"
        
        topics_to_cover = module.get("topics", [module["title"]])
        slides_per_topic = max(2, (max_slides - 1) // max(1, len(topics_to_cover)))
        
        # Generate slides for each unique topic
        for idx, topic in enumerate(topics_to_cover):
            if topic in covered_topics:
                continue
            
            # Add delay between API calls to avoid rate limits
            if idx > 0:
                await asyncio.sleep(2)  # 2 second delay between topics
                logger.info(f"Rate limit delay: waiting 2 seconds before next topic...")
            
            try:
                logger.info(f"Processing topic {idx+1}/{len(topics_to_cover)}: {topic}")
                
                topic_slides = await self.generate_slides_batch(
                    topic, 
                    module, 
                    num_slides=slides_per_topic,
                    previous_context=context_history
                )
                
                # Only add validated slides
                if topic_slides:
                    for slide in topic_slides:
                        if len(slides) >= max_slides:
                            break
                        
                        slides.append(slide)
                        
                        # Update context history with detailed info
                        slide_summary = f"\nSlide: {slide.get('title', 'Untitled')}\nContent: {' '.join(slide.get('content', []))[:200]}\n"
                        context_history += slide_summary
                        covered_topics.add(topic)
                else:
                    logger.warning(f"No slides generated for topic: {topic}")
                    # Create fallback slide
                    if len(slides) < max_slides:
                        fallback = self._create_fallback_slide(topic, module)
                        slides.append(fallback)
                        covered_topics.add(topic)
                
            except Exception as e:
                logger.error(f"Error generating slides for {topic}: {e}")
                # Fallback for errors
                if len(slides) < max_slides:
                    fallback = self._create_fallback_slide(topic, module)
                    slides.append(fallback)
                    covered_topics.add(topic)
        
        logger.info(f"Module {module['title']}: Generated {len(slides)} slides")
        return slides[:max_slides]
    
    async def generate_presentation(self, outline: Dict, max_slides_per_module: int = 12, include_assessments: bool = True) -> Tuple[List[Dict], Dict]:
        all_slides = []
        all_assessments = {}
        
        description = outline.get("description", "Comprehensive professional training")
        title_slide = {
            "slide_type": "title",
            "title": outline.get("title", "Professional Training Program"),
            "content": description,
            "notes": f"Welcome to {outline.get('title')}. {description} This training is designed for {outline.get('target_audience', 'professionals')}.",
            "estimated_time_minutes": 3
        }
        all_slides.append(title_slide)
        
        agenda_slide = {
            "slide_type": "agenda",
            "title": "Training Agenda",
            "content": [f"{i+1}. {mod['title']} ({mod.get('duration', '45 min')})" for i, mod in enumerate(outline.get("modules", []))],
            "notes": f"Today we'll cover {len(outline.get('modules', []))} key modules. {' '.join([m['title'] for m in outline.get('modules', [])[:2]])}",
            "estimated_time_minutes": 2
        }
        all_slides.append(agenda_slide)
        
        objectives_slide = {
            "slide_type": "objectives",
            "title": "Learning Objectives",
            "content": outline.get("objectives", ["Master key concepts"]),
            "notes": f"By the end of this training, you will: {' '.join(outline.get('objectives', [])[:2])}",
            "estimated_time_minutes": 3
        }
        all_slides.append(objectives_slide)
        
        modules = outline.get("modules", [])
        total_modules = len(modules)
        
        # Recalculate slides per module based on actual module count
        # Reserve 3 slides for intro/outro
        content_slides_needed = max_slides_per_module * total_modules
        adjusted_slides_per_module = max(8, content_slides_needed // total_modules)
        
        logger.info(f"Target: {content_slides_needed} content slides across {total_modules} modules")
        logger.info(f"Adjusted to {adjusted_slides_per_module} slides per module")
        
        for module_idx, module in enumerate(modules):
            try:
                logger.info(f"Generating slides for module {module_idx + 1}/{total_modules}: {module.get('title')}")
                
                # Add delay between modules to avoid rate limits
                if module_idx > 0:
                    delay = 3
                    logger.info(f"Rate limit delay: waiting {delay} seconds before next module...")
                    await asyncio.sleep(delay)
                
                module_slides = await self.generate_slides_for_module(module, adjusted_slides_per_module)
                
                if module_slides:
                    all_slides.extend(module_slides)
                    logger.info(f"Successfully generated {len(module_slides)} slides for module {module_idx + 1}")
                else:
                    logger.warning(f"No slides generated for module {module_idx + 1}, adding placeholder")
                    # Add placeholder slide
                    all_slides.append({
                        "slide_type": "title",
                        "title": module.get("title", "Module"),
                        "content": module.get("objectives", []),
                        "notes": " ".join(module.get("key_points", [])),
                        "estimated_time_minutes": 5
                    })
                
                # Generate assessment for this module if enabled
                if include_assessments:
                    try:
                        await asyncio.sleep(2)  # Delay before assessment generation
                        assessment = await self.assessment_generator.generate_assessment(module, num_questions=5)
                        all_assessments[module['title']] = assessment
                        logger.info(f"Generated assessment for module {module_idx + 1}")
                    except Exception as e:
                        logger.error(f"Error generating assessment for {module['title']}: {e}")
                
            except Exception as e:
                logger.error(f"Error processing module {module.get('title')}: {e}")
                logger.error(f"Error details: {type(e).__name__}: {str(e)}")
                # Add minimal placeholder slide
                all_slides.append({
                    "slide_type": "title",
                    "title": module.get("title", "Module"),
                    "content": module.get("objectives", ["Module content"]),
                    "notes": " ".join(module.get("key_points", ["Module information"])),
                    "estimated_time_minutes": 5
                })
        
        module_titles = [m['title'] for m in outline.get('modules', [])]
        closing_slide = {
            "slide_type": "closing",
            "title": "Summary & Next Steps",
            "content": [
                f"Mastered: {', '.join(module_titles[:2])}" if len(module_titles) >= 2 else "Key concepts reviewed",
                "Apply this knowledge immediately in your role",
                "Complete the final assessment" if include_assessments else "Review key materials",
                "Questions and discussion"
            ],
            "notes": f"Today we covered {len(module_titles)} modules: {', '.join(module_titles)}. Remember to apply these concepts in your daily work.",
            "estimated_time_minutes": 5
        }
        all_slides.append(closing_slide)
        
        logger.info(f"Presentation complete: {len(all_slides)} total slides generated")
        return all_slides, all_assessments
    
    def _apply_brand_colors(self, prs: PptxPresentation):
        """Apply brand colors to presentation if configured"""
        if not self.brand_config or 'colors' not in self.brand_config:
            return
        
        # This is a simplified version - in production you'd apply colors to master slides
        pass
    
    def create_powerpoint(self, slides: List[Dict], output_path: Union[str, Path]) -> Path:
        if isinstance(output_path, str):
            output_path = Path(output_path)
        
        prs = PptxPresentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        self._apply_brand_colors(prs)
        
        for slide_data in slides:
            slide_type = slide_data.get("slide_type", "content")
            
            try:
                if slide_type in ["title", "closing"]:
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    title = slide.shapes.title
                    subtitle = slide.placeholders[1]
                    
                    title.text = slide_data["title"]
                    
                    if isinstance(slide_data["content"], list):
                        subtitle.text = "\n".join(str(item) for item in slide_data["content"])
                    else:
                        subtitle.text = str(slide_data["content"])
                
                else:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    
                    title.text = slide_data["title"]
                    
                    text_frame = content.text_frame
                    text_frame.clear()
                    
                    bullet_points = slide_data["content"] if isinstance(slide_data["content"], list) else [str(slide_data["content"])]
                    
                    for bullet in bullet_points:
                        p = text_frame.add_paragraph()
                        p.text = str(bullet)
                        p.level = 0
                
                if "notes" in slide_data:
                    notes_slide = slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    text_frame.text = slide_data["notes"]
                    
            except Exception as e:
                logger.error(f"Error creating slide '{slide_data.get('title', 'Unknown')}': {e}")
        
        try:
            prs.save(output_path)
            logger.info(f"Successfully saved presentation to: {output_path}")
            return output_path
        except Exception as e:
            logger.error(f"Error saving presentation: {e}")
            raise
    
    async def create_assessment_document(self, assessments: Dict, outline: Dict, output_path: Union[str, Path]) -> Path:
        """Create a document with all assessments"""
        if isinstance(output_path, str):
            output_path = Path(output_path)
        
        doc = DocxDocument()
        
        doc.add_heading(f"Assessment: {outline.get('title', 'Training')}", 0)
        
        doc.add_paragraph(f"Training Program: {outline.get('title')}")
        doc.add_paragraph(f"Duration: {outline.get('estimated_duration_hours', 8)} hours")
        doc.add_paragraph("")
        
        for module_title, assessment in assessments.items():
            doc.add_heading(f"Module: {module_title}", 1)
            
            questions = assessment.get("questions", [])
            
            for idx, question in enumerate(questions, 1):
                q_type = question.get("type", "multiple_choice")
                
                doc.add_heading(f"Question {idx} ({q_type.replace('_', ' ').title()})", 2)
                
                if q_type == "multiple_choice":
                    doc.add_paragraph(question.get("question", ""))
                    for opt in question.get("options", []):
                        doc.add_paragraph(f"  {opt}", style='List Bullet')
                    doc.add_paragraph(f"Correct Answer: {question.get('correct_answer')}")
                    doc.add_paragraph(f"Explanation: {question.get('explanation')}")
                
                elif q_type == "true_false":
                    doc.add_paragraph(question.get("question", ""))
                    doc.add_paragraph(f"Correct Answer: {question.get('correct_answer')}")
                    doc.add_paragraph(f"Explanation: {question.get('explanation')}")
                
                elif q_type == "scenario":
                    doc.add_paragraph("Scenario:")
                    doc.add_paragraph(question.get("scenario", ""))
                    doc.add_paragraph(f"Question: {question.get('question')}")
                    doc.add_paragraph("Suggested Answer:")
                    doc.add_paragraph(question.get("suggested_answer", ""))
                    doc.add_paragraph("Grading Rubric:")
                    for criterion in question.get("rubric", []):
                        doc.add_paragraph(f"• {criterion}")
                
                doc.add_paragraph("")
        
        doc.save(output_path)
        logger.info(f"Successfully saved assessment document to: {output_path}")
        return output_path
    
    @retry(wait=wait_exponential(multiplier=1, min=4, max=60), stop=stop_after_attempt(3))
    async def generate_trainer_guide_section(self, module: Dict, slide: Dict) -> Dict:
        prompt = PROMPTS["trainer_guide"].format(
            module=json.dumps(module, indent=2),
            slide=json.dumps(slide, indent=2)
        )
        
        response = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are creating a comprehensive trainer guide. Write the ACTUAL CONTENT the trainer should teach - specific facts, procedures, examples, and explanations with real information from the source. Do NOT write meta-instructions. Minimum 300 words of teaching content per section."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            response_format={"type": "json_object"},
            max_tokens=2000  # Increased for more detailed content
        )
        
        result = response.choices[0].message.content
        return json.loads(result)
    
    async def generate_trainer_guide(self, outline: Dict, slides: List[Dict], assessments: Dict, output_path: Union[str, Path]) -> Path:
        if isinstance(output_path, str):
            output_path = Path(output_path)
        
        doc = DocxDocument()
        
        doc.add_heading(f"Trainer Guide: {outline.get('title', 'Training Program')}", 0)
        
        # Calculate total duration
        total_slide_time = sum(slide.get("estimated_time_minutes", 5) for slide in slides)
        duration_info = calculate_training_duration(len(slides), "medium", True)
        
        doc.add_heading("Training Overview", 1)
        doc.add_paragraph(outline.get("description", "Comprehensive training program"))
        doc.add_paragraph(f"Estimated Duration: {format_duration_display(duration_info)}")
        doc.add_paragraph(f"Total Slides: {len(slides)}")
        
        doc.add_heading("Target Audience", 2)
        doc.add_paragraph(outline.get("target_audience", "Professionals"))
        
        doc.add_heading("Prerequisites", 2)
        doc.add_paragraph(outline.get("prerequisites", "Basic knowledge"))
        
        doc.add_heading("Learning Objectives", 2)
        for objective in outline.get("objectives", []):
            doc.add_paragraph(f"• {objective}")
        
        doc.add_heading("Training Agenda", 1)
        for i, module in enumerate(outline.get("modules", []), 1):
            doc.add_paragraph(f"{i}. {module.get('title')} - {module.get('duration', '45 min')}")
        
        current_module = None
        
        for slide in slides:
            try:
                if slide.get("slide_type") == "title" and slide.get("title") in [m.get("title") for m in outline.get("modules", [])]:
                    for module in outline.get("modules", []):
                        if module.get("title") == slide.get("title"):
                            current_module = module
                            
                            doc.add_page_break()
                            doc.add_heading(f"Module: {module.get('title')}", 1)
                            
                            doc.add_heading("Module Objectives", 2)
                            for obj in module.get("objectives", []):
                                doc.add_paragraph(f"• {obj}")
                            
                            doc.add_heading("Key Concepts", 2)
                            for point in module.get("key_points", []):
                                doc.add_paragraph(f"• {point}")
                            
                            # Add assessment info if available
                            if module.get('title') in assessments:
                                doc.add_heading("Assessment", 2)
                                doc.add_paragraph(f"This module includes {len(assessments[module.get('title')].get('questions', []))} assessment questions.")
                            
                            break
                
                if slide.get("slide_type") in ["title", "agenda", "objectives"] and slides.index(slide) < 3:
                    continue
                
                if current_module:
                    guide_section = await self.generate_trainer_guide_section(current_module, slide)
                    
                    doc.add_heading(f"Slide: {slide.get('title')}", 2)
                    
                    # Teaching Content - The main substance
                    doc.add_heading("Teaching Content", 3)
                    doc.add_paragraph(guide_section.get("teaching_content", "Detailed teaching content"))
                    
                    doc.add_heading("Delivery Approach", 3)
                    doc.add_paragraph(guide_section.get("delivery_approach", "Teaching methods"))
                    
                    # Real Examples
                    doc.add_heading("Examples from Source Material", 3)
                    for example in guide_section.get("real_examples", []):
                        doc.add_paragraph(f"• {example}")
                    
                    doc.add_heading("Key Points to Emphasize", 3)
                    for point in guide_section.get("key_points_to_emphasize", []):
                        doc.add_paragraph(f"• {point}")
                    
                    # Interactive Activity
                    if "interactive_activity" in guide_section:
                        activity = guide_section["interactive_activity"]
                        doc.add_heading("Interactive Activity", 3)
                        doc.add_paragraph(f"Type: {activity.get('activity_type', 'Discussion')}")
                        doc.add_paragraph(f"Instructions: {activity.get('instructions', '')}")
                        doc.add_paragraph(f"Duration: {activity.get('duration', '5 minutes')}")
                        if activity.get('materials'):
                            doc.add_paragraph("Materials needed:")
                            for material in activity['materials']:
                                doc.add_paragraph(f"  • {material}")
                    
                    # Q&A Section
                    doc.add_heading("Anticipated Questions & Answers", 3)
                    for qa in guide_section.get("anticipated_questions", []):
                        doc.add_paragraph(f"Q: {qa.get('question', '')}")
                        doc.add_paragraph(f"A: {qa.get('detailed_answer', '')}")
                        doc.add_paragraph("")
                    
                    doc.add_heading("Assessment Checkpoint", 3)
                    doc.add_paragraph(guide_section.get("assessment_checkpoint", "Quick knowledge check"))
                    
                    doc.add_heading("Timing & Transitions", 3)
                    doc.add_paragraph(f"Time: {guide_section.get('timing_guidance', '5-10 minutes')}")
                    doc.add_paragraph(f"Transitions: {guide_section.get('transitions', 'Smooth transition')}")
                    
                    doc.add_paragraph("-" * 50)
                    
            except Exception as e:
                logger.error(f"Error in trainer guide section for slide '{slide.get('title')}': {e}")
        
        try:
            doc.save(output_path)
            logger.info(f"Successfully saved trainer guide to: {output_path}")
            return output_path
        except Exception as e:
            logger.error(f"Error saving trainer guide: {e}")
            raise

class ProgressTracker:
    def __init__(self):
        self.status = "idle"
        self.step = ""
        self.progress = 0
        self.max_progress = 100
        self.message = ""
        self.start_time = None
        self.error = None
    
    def start(self, max_steps: int = 100):
        self.status = "running"
        self.step = "initializing"
        self.progress = 0
        self.max_progress = max_steps
        self.message = "Starting process..."
        self.start_time = time.time()
        self.error = None
    
    def update(self, step: str, progress: int, message: str = ""):
        self.step = step
        self.progress = min(progress, self.max_progress)
        self.message = message or f"Processing {step}..."
    
    def complete(self, message: str = "Process completed successfully"):
        self.status = "complete"
        self.progress = self.max_progress
        self.message = message
    
    def fail(self, error: str):
        self.status = "failed"
        self.error = error
        self.message = f"Error: {error}"
    
    def get_info(self) -> Dict:
        elapsed = 0
        if self.start_time:
            elapsed = time.time() - self.start_time
        return {
            "status": self.status,
            "step": self.step,
            "progress": self.progress,
            "max_progress": self.max_progress,
            "percent": int(100 * self.progress / max(1, self.max_progress)),
            "message": self.message,
            "elapsed": elapsed,
            "error": self.error
        }

class TrainingGenerator:
    def __init__(self, openai_api_key: str, brand_config: Optional[Dict] = None):
        self.document_processor = DocumentProcessor()
        self.content_analyzer = ContentAnalyzer(openai_api_key)
        self.document_generator = DocumentGenerator(openai_api_key, brand_config)
        self.progress_tracker = ProgressTracker()
    
    def set_brand_config(self, brand_config: Dict):
        self.document_generator.set_brand_config(brand_config)
    
    async def process_document(self, file_path: Union[str, Path]) -> str:
        self.progress_tracker.update("document_processing", 10, "Processing document...")
        content = self.document_processor.process_document(file_path)
        self.document_generator.set_source_content(content)
        return content
    
    async def analyze_content(self, content: str, duration: str = "1 day") -> Tuple[List[Dict], Dict]:
        self.progress_tracker.update("topic_extraction", 20, "Extracting topics...")
        topics = self.content_analyzer.extract_topics(content)
        
        self.progress_tracker.update("outline_generation", 30, "Generating outline...")
        outline = self.content_analyzer.generate_outline(topics, content, duration)
        
        return topics, outline
    
    async def analyze_quality(self, content: str, outline: Dict) -> Dict:
        self.progress_tracker.update("quality_analysis", 35, "Analyzing content quality...")
        return self.content_analyzer.analyze_quality(content, outline)
    
    async def generate_materials(self, outline: Dict, output_dir: Union[str, Path], 
                                 max_slides_per_module: int, include_assessments: bool = True) -> Dict:
        if isinstance(output_dir, str):
            output_dir = Path(output_dir)
        
        output_dir.mkdir(exist_ok=True, parents=True)
        
        self.progress_tracker.update("slide_generation", 40, "Generating slides...")
        slides, assessments = await self.document_generator.generate_presentation(
            outline, max_slides_per_module, include_assessments
        )
        
        # Calculate duration
        duration_info = calculate_training_duration(len(slides), "medium", include_assessments)
        
        self.progress_tracker.update("presentation_creation", 60, "Creating PowerPoint...")
        safe_title = outline.get('title', 'Training').replace(' ', '_').replace('/', '_')
        pptx_path = output_dir / f"{safe_title}.pptx"
        pptx_path = self.document_generator.create_powerpoint(slides, pptx_path)
        
        self.progress_tracker.update("trainer_guide_creation", 75, "Creating trainer guide...")
        guide_path = output_dir / f"{safe_title}_TrainerGuide.docx"
        guide_path = await self.document_generator.generate_trainer_guide(outline, slides, assessments, guide_path)
        
        assessment_path = None
        if include_assessments and assessments:
            self.progress_tracker.update("assessment_creation", 85, "Creating assessment document...")
            assessment_path = output_dir / f"{safe_title}_Assessment.docx"
            assessment_path = await self.document_generator.create_assessment_document(assessments, outline, assessment_path)
        
        if not pptx_path.exists() or not guide_path.exists():
            raise FileNotFoundError("Output files were not created")
        
        return {
            "pptx_path": str(pptx_path),
            "guide_path": str(guide_path),
            "assessment_path": str(assessment_path) if assessment_path else None,
            "duration_info": duration_info,
            "total_slides": len(slides)
        }
    
    async def generate_training_materials(self, file_path: Union[str, Path], output_dir: Union[str, Path], 
                                         max_slides_per_module: int, include_assessments: bool = True,
                                         duration: str = "1 day") -> Dict:
        try:
            self.progress_tracker.start(100)
            
            content = await self.process_document(file_path)
            
            if not content or len(content.strip()) < 100:
                raise ValueError("Insufficient content extracted from document")
            
            topics, outline = await self.analyze_content(content, duration)
            quality_analysis = await self.analyze_quality(content, outline)
            
            result = await self.generate_materials(outline, output_dir, max_slides_per_module, include_assessments)
            result["topics"] = topics
            result["outline"] = outline
            result["quality_analysis"] = quality_analysis
            
            self.progress_tracker.complete("Training materials generated successfully")
            
            return result
            
        except Exception as e:
            logger.error(f"Error generating training materials: {e}")
            self.progress_tracker.fail(str(e))
            raise

# Custom CSS for beautiful UI
def load_custom_css():
    if not OPENAI_API_KEY:
        st.error("⚠️ API Key Not Configured")
        st.markdown("""
        ### How to fix this:
        
        **For Streamlit Cloud:**
        1. Go to your app settings → **Secrets**
        2. Add the following:
        ```toml
        OPENAI_API_KEY = "sk-proj-your-api-key-here"
        ```
        3. Click **Save** and wait for the app to redeploy
        
        **For local development:**
        - Set environment variable: `export OPENAI_API_KEY="your-key"`
        """)
        st.stop()
    st.markdown("""
    <style>
    /* Main theme colors */
    :root {
        --primary-color: #4F46E5;
        --secondary-color: #10B981;
        --background-color: #F9FAFB;
        --card-background: #FFFFFF;
        --text-color: #1F2937;
        --border-color: #E5E7EB;
    }
    
    /* Main container styling */
    .main {
        background-color: var(--background-color);
    }
    
    /* Card styling */
    .custom-card {
        background-color: var(--card-background);
        padding: 2rem;
        border-radius: 1rem;
        box-shadow: 0 1px 3px 0 rgba(0, 0, 0, 0.1);
        margin-bottom: 1.5rem;
        border: 1px solid var(--border-color);
    }
    
    /* File upload area */
    .upload-area {
        border: 2px dashed var(--border-color);
        border-radius: 0.5rem;
        padding: 2rem;
        text-align: center;
        background-color: var(--background-color);
        transition: all 0.3s ease;
    }
    
    .upload-area:hover {
        border-color: var(--primary-color);
        background-color: #EEF2FF;
    }
    
    /* Button styling */
    .stButton>button {
        background-color: var(--primary-color);
        color: white;
        border: none;
        padding: 0.75rem 2rem;
        border-radius: 0.5rem;
        font-weight: 500;
        transition: all 0.3s ease;
    }
    
    .stButton>button:hover {
        background-color: #4338CA;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1);
    }
    
    /* Topic card styling */
    .topic-card {
        background: white;
        border: 1px solid var(--border-color);
        border-radius: 0.5rem;
        padding: 1rem;
        margin-bottom: 1rem;
        transition: all 0.2s ease;
    }
    
    .topic-card:hover {
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1);
        border-color: var(--primary-color);
    }
    
    .topic-badge {
        display: inline-block;
        padding: 0.25rem 0.75rem;
        border-radius: 9999px;
        font-size: 0.75rem;
        font-weight: 500;
        margin-left: 0.5rem;
    }
    
    .badge-high {
        background-color: #FEE2E2;
        color: #991B1B;
    }
    
    .badge-medium {
        background-color: #FEF3C7;
        color: #92400E;
    }
    
    .badge-low {
        background-color: #DBEAFE;
        color: #1E40AF;
    }
    
    /* Module editor styling */
    .module-editor {
        background: white;
        border: 1px solid var(--border-color);
        border-radius: 0.5rem;
        padding: 1.5rem;
        margin-bottom: 1rem;
    }
    
    .module-header {
        display: flex;
        justify-content: space-between;
        align-items: center;
        margin-bottom: 1rem;
        padding-bottom: 1rem;
        border-bottom: 1px solid var(--border-color);
    }
    
    /* Progress bar styling */
    .stProgress > div > div {
        background-color: var(--primary-color);
    }
    
    /* Info boxes */
    .info-box {
        padding: 1rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
        border-left: 4px solid;
    }
    
    .info-box.success {
        background-color: #D1FAE5;
        border-color: var(--secondary-color);
        color: #065F46;
    }
    
    .info-box.info {
        background-color: #DBEAFE;
        border-color: #3B82F6;
        color: #1E40AF;
    }
    
    .info-box.warning {
        background-color: #FEF3C7;
        border-color: #F59E0B;
        color: #92400E;
    }
    
    /* Metric card */
    .metric-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 0.5rem;
        text-align: center;
        margin-bottom: 1rem;
    }
    
    .metric-value {
        font-size: 2rem;
        font-weight: 700;
        margin: 0.5rem 0;
    }
    
    .metric-label {
        font-size: 0.875rem;
        opacity: 0.9;
    }
    
    /* Hide Streamlit branding */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    
    /* Headers */
    h1 {
        color: var(--text-color);
        font-weight: 700;
    }
    
    h2 {
        color: var(--text-color);
        font-weight: 600;
        margin-top: 1.5rem;
    }
    
    h3 {
        color: var(--text-color);
        font-weight: 500;
    }
    </style>
    """, unsafe_allow_html=True)

def render_step_indicator(current_step: str):
    steps = {
        "upload": {"number": 1, "title": "Upload Documents"},
        "topics": {"number": 2, "title": "Review Topics"},
        "edit": {"number": 3, "title": "Edit Outline"},
        "generate": {"number": 4, "title": "Generate Materials"}
    }
    
    step_order = ["upload", "topics", "edit", "generate"]
    current_index = step_order.index(current_step)
    
    cols = st.columns(4)
    
    for i, (col, step_key) in enumerate(zip(cols, step_order)):
        step = steps[step_key]
        with col:
            if i < current_index:
                # Completed step
                st.markdown(f"""
                <div style="text-align: center; padding: 1rem;">
                    <div style="width: 40px; height: 40px; border-radius: 50%; background-color: #10B981; color: white; 
                                display: inline-flex; align-items: center; justify-content: center; font-weight: 600; 
                                margin: 0 auto 0.5rem;">
                        {step["number"]}
                    </div>
                    <div style="font-size: 0.75rem; color: #1F2937; font-weight: 500;">
                        {step["title"]}
                    </div>
                </div>
                """, unsafe_allow_html=True)
            elif i == current_index:
                # Active step
                st.markdown(f"""
                <div style="text-align: center; padding: 1rem;">
                    <div style="width: 40px; height: 40px; border-radius: 50%; background-color: #4F46E5; color: white; 
                                display: inline-flex; align-items: center; justify-content: center; font-weight: 600; 
                                margin: 0 auto 0.5rem; box-shadow: 0 0 0 4px rgba(79, 70, 229, 0.2);">
                        {step["number"]}
                    </div>
                    <div style="font-size: 0.75rem; color: #4F46E5; font-weight: 600;">
                        {step["title"]}
                    </div>
                </div>
                """, unsafe_allow_html=True)
            else:
                # Upcoming step
                st.markdown(f"""
                <div style="text-align: center; padding: 1rem;">
                    <div style="width: 40px; height: 40px; border-radius: 50%; background-color: #E5E7EB; color: #6B7280; 
                                display: inline-flex; align-items: center; justify-content: center; font-weight: 600; 
                                margin: 0 auto 0.5rem;">
                        {step["number"]}
                    </div>
                    <div style="font-size: 0.75rem; color: #6B7280; font-weight: 500;">
                        {step["title"]}
                    </div>
                </div>
                """, unsafe_allow_html=True)

def render_file_card(filename: str, filesize: str, index: int):
    st.markdown(f"""
    <div class="topic-card">
        <div style="display: flex; justify-content: space-between; align-items: center;">
            <div>
                <strong>{filename}</strong>
                <div style="font-size: 0.875rem; color: #6B7280; margin-top: 0.25rem;">
                    Size: {filesize}
                </div>
            </div>
            <div style="color: var(--secondary-color); font-weight: 500;">
                Uploaded
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)

def main():
    st.set_page_config(
        page_title="QUINNS Training Generator Pro",
        page_icon="📚",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    load_custom_css()
    
    # Initialize session state
    if 'generator' not in st.session_state:
        st.session_state.generator = None
    if 'api_key' not in st.session_state:
        st.session_state.api_key = ""
    if 'topics' not in st.session_state:
        st.session_state.topics = []
    if 'outline' not in st.session_state:
        st.session_state.outline = {}
    if 'content' not in st.session_state:
        st.session_state.content = ""
    if 'step' not in st.session_state:
        st.session_state.step = "upload"
    if 'output_files' not in st.session_state:
        st.session_state.output_files = {}
    if 'file_paths' not in st.session_state:
        st.session_state.file_paths = []
    if 'max_slides_per_module' not in st.session_state:
        st.session_state.max_slides_per_module = 12
    if 'target_total_slides' not in st.session_state:
        st.session_state.target_total_slides = 36
    if 'combined_content' not in st.session_state:
        st.session_state.combined_content = ""
    if 'template' not in st.session_state:
        st.session_state.template = "corporate"
    if 'training_duration' not in st.session_state:
        st.session_state.training_duration = "1 day"
    if 'include_assessments' not in st.session_state:
        st.session_state.include_assessments = True
    if 'pace' not in st.session_state:
        st.session_state.pace = "medium"
    if 'quality_analysis' not in st.session_state:
        st.session_state.quality_analysis = {}
    if 'brand_config' not in st.session_state:
        st.session_state.brand_config = {}
    
    # Initialize additional session state for sidebar
    if 'training_type' not in st.session_state:
        st.session_state.training_type = "Auto-Detect"
    if 'session_duration' not in st.session_state:
        st.session_state.session_duration = "Half Day (4 hours)"
    if 'target_slide_count' not in st.session_state:
        st.session_state.target_slide_count = "Auto-calculated"
    if 'include_activities' not in st.session_state:
        st.session_state.include_activities = True
    if 'add_break_slides' not in st.session_state:
        st.session_state.add_break_slides = False
    if 'enhanced_speaker_notes' not in st.session_state:
        st.session_state.enhanced_speaker_notes = True
    if 'professional_templates' not in st.session_state:
        st.session_state.professional_templates = True
    
    # Set API key from environment or hardcoded value (for backend configuration)
    if not st.session_state.generator:
        try:
            st.session_state.generator = TrainingGenerator(
            OPENAI_API_KEY,  # Use module-level variable
            st.session_state.brand_config
        )
        except Exception as e:
            st.error(f"Error: {e}")
            st.stop()
    
    # Sidebar configuration with modern design
    with st.sidebar:
        st.markdown("""
        <style>
        /* Sidebar styling with light background and dark text */
        [data-testid="stSidebar"] {
            background-color: #F8F9FA;
        }
        
        [data-testid="stSidebar"] .element-container {
            color: #1F2937;
        }
        
        /* Section headers */
        .sidebar-section-header {
            color: #1F2937;
            font-size: 1.25rem;
            font-weight: 700;
            margin: 1.5rem 0 1rem 0;
            display: flex;
            align-items: center;
            gap: 0.5rem;
            padding: 0.75rem;
            background: white;
            border-radius: 0.5rem;
            border-left: 4px solid #667eea;
        }
        
        /* Configuration labels */
        .config-field-label {
            color: #374151;
            font-size: 0.875rem;
            font-weight: 600;
            margin-bottom: 0.5rem;
            margin-top: 0.75rem;
        }
        
        /* Enhancement section */
        .enhancement-header {
            color: #1F2937;
            font-size: 1.25rem;
            font-weight: 700;
            margin: 1.5rem 0 1rem 0;
            display: flex;
            align-items: center;
            gap: 0.5rem;
            padding: 0.75rem;
            background: white;
            border-radius: 0.5rem;
            border-left: 4px solid #F59E0B;
        }
        
        /* Divider */
        .sidebar-divider {
            border-top: 2px solid #E5E7EB;
            margin: 1.5rem 0;
        }
        
        /* Help icon styling */
        .help-icon {
            color: #6B7280;
            cursor: help;
        }
        
        /* Toggle labels */
        [data-testid="stSidebar"] label {
            color: #1F2937 !important;
            font-weight: 500;
        }
        
        /* Selectbox labels */
        [data-testid="stSidebar"] .stSelectbox label {
            color: #374151 !important;
            font-weight: 600;
        }
        </style>
        """, unsafe_allow_html=True)
        
        # Training Configuration Section
        st.markdown('<div class="sidebar-section-header">🎯 Training Configuration</div>', unsafe_allow_html=True)
        
        # Training Type
        st.markdown('<div class="config-field-label">Training Type</div>', unsafe_allow_html=True)
        training_types = [
            "Auto-Detect",
            "Corporate Training",
            "Technical Skills",
            "Compliance & Safety",
            "New Employee Onboarding",
            "Leadership Development",
            "Sales Training",
            "Customer Service",
            "Soft Skills",
            "Product Training"
        ]
        st.session_state.training_type = st.selectbox(
            "Training Type",
            options=training_types,
            index=training_types.index(st.session_state.training_type),
            label_visibility="collapsed",
            help="Select training type or let AI auto-detect from your content"
        )
        
        # Session Duration
        st.markdown('<div class="config-field-label">Session Duration</div>', unsafe_allow_html=True)
        duration_options = [
            "30 Minutes",
            "1 Hour",
            "2 Hours",
            "Half Day (4 hours)",
            "Full Day (8 hours)",
            "2 Days",
            "3 Days",
            "1 Week"
        ]
        st.session_state.session_duration = st.selectbox(
            "Session Duration",
            options=duration_options,
            index=duration_options.index(st.session_state.session_duration),
            label_visibility="collapsed",
            help="Target duration for your training program"
        )
        
        # Map duration to internal format
        duration_mapping = {
            "30 Minutes": "30 minutes",
            "1 Hour": "1 hour",
            "2 Hours": "2 hours",
            "Half Day (4 hours)": "half day",
            "Full Day (8 hours)": "1 day",
            "2 Days": "2 days",
            "3 Days": "3 days",
            "1 Week": "1 week"
        }
        st.session_state.training_duration = duration_mapping[st.session_state.session_duration]
        
        # Target Slide Count
        st.markdown('<div class="config-field-label">Target Slide Count</div>', unsafe_allow_html=True)
        slide_count_option = st.selectbox(
            "Target Slide Count",
            options=["Auto-calculated based on duration", "10-15 slides", "20-30 slides", "30-50 slides", "50+ slides"],
            label_visibility="collapsed",
            help="Number of slides to generate (auto-calculated recommended)"
        )
        
        # Map slide count to slides per module
        if "10-15" in slide_count_option:
            st.session_state.max_slides_per_module = 5
        elif "20-30" in slide_count_option:
            st.session_state.max_slides_per_module = 8
        elif "30-50" in slide_count_option:
            st.session_state.max_slides_per_module = 12
        elif "50+" in slide_count_option:
            st.session_state.max_slides_per_module = 20
        else:
            # Auto-calculate based on duration - FIXED TO MATCH TARGET DURATION
            # Calculate slides needed to fill the target duration
            # Assume 5 minutes per slide average
            duration_to_minutes = {
                "30 minutes": 30,
                "1 hour": 60,
                "2 hours": 120,
                "half day": 240,  # 4 hours
                "1 day": 480,     # 8 hours
                "2 days": 960,    # 16 hours
                "3 days": 1440,   # 24 hours
                "1 week": 2400    # 40 hours (5 days)
            }
            
            target_minutes = duration_to_minutes.get(st.session_state.training_duration, 240)
            
            # Calculate total slides needed (5 min per slide + 20% for activities + breaks)
            slides_needed = int(target_minutes / 5 * 0.75)  # 75% efficiency for activities/breaks
            
            # Estimate number of modules (will be adjusted after outline generation)
            estimated_modules = min(5, max(3, slides_needed // 15))  # 3-5 modules typical
            
            # Calculate slides per module
            st.session_state.max_slides_per_module = max(10, slides_needed // estimated_modules)
            
            # Store total target slides for later use
            st.session_state.target_total_slides = slides_needed
        
        st.markdown('<div class="sidebar-divider"></div>', unsafe_allow_html=True)
        
        # Enhancement Options Section
        st.markdown('<div class="enhancement-header">✨ Enhancement Options</div>', unsafe_allow_html=True)
        
        st.session_state.include_activities = st.toggle(
            "Include Interactive Activities",
            value=st.session_state.include_activities,
            help="Add hands-on exercises, group discussions, and interactive elements"
        )
        
        st.session_state.add_break_slides = st.toggle(
            "Add Break Slides",
            value=st.session_state.add_break_slides,
            help="Insert break reminder slides at regular intervals"
        )
        
        st.session_state.enhanced_speaker_notes = st.toggle(
            "Enhanced Speaker Notes",
            value=st.session_state.enhanced_speaker_notes,
            help="Generate detailed presenter notes with timing and tips"
        )
        
        st.session_state.professional_templates = st.toggle(
            "Professional Templates",
            value=st.session_state.professional_templates,
            help="Use professionally designed slide layouts and themes"
        )
        
        # Update include_assessments from enhanced options
        st.session_state.include_assessments = st.toggle(
            "Generate Assessments",
            value=st.session_state.get('include_assessments', True),
            help="Create quiz questions and knowledge checks"
        )
        
        st.markdown('<div class="sidebar-divider"></div>', unsafe_allow_html=True)
        
        # Show estimated generation target
        if 'target_total_slides' in st.session_state and st.session_state.target_total_slides:
            st.markdown("### 📈 Generation Target")
            st.info(f"**Target Slides:** {st.session_state.target_total_slides} slides\n\n"
                   f"**Per Module:** ~{st.session_state.max_slides_per_module} slides\n\n"
                   f"**Estimated Duration:** {st.session_state.session_duration}")
        
        st.markdown('<div class="sidebar-divider"></div>', unsafe_allow_html=True)
        
        # Quick Stats
        st.markdown("### 📊 Quick Stats")
        if st.session_state.step != "upload":
            if st.session_state.file_paths:
                st.metric("Documents Uploaded", len(st.session_state.file_paths))
            
            if st.session_state.topics:
                st.metric("Topics Identified", len(st.session_state.topics))
            
            if st.session_state.outline and st.session_state.outline.get('modules'):
                st.metric("Modules Created", len(st.session_state.outline['modules']))
        
        st.markdown('<div class="sidebar-divider"></div>', unsafe_allow_html=True)
        
        # About Section
        with st.expander("ℹ️ About QUINNS Pro"):
            st.markdown("""
            **Advanced Features:**
            - AI-powered content analysis
            - Auto-detected training types
            - Smart slide calculation
            - Quality metrics
            - Assessment generation
            - Brand customization
            - Multiple export formats
            """)
        
        # Tips Section
        with st.expander("💡 Tips & Best Practices"):
            st.markdown("""
            **For Best Results:**
            1. Upload complete, well-structured documents
            2. Let AI auto-detect training type
            3. Use auto-calculated slide count
            4. Enable interactive activities
            5. Review quality metrics
            6. Customize outline before generation
            """)
    
    # Main content area
    LOGO_IMAGE_PATH = "logo_url.png"
    st.image(
                LOGO_IMAGE_PATH, 
                caption=None, # No caption for a logo
                width=150 # Adjust width as needed for a logo
            )
    st.title("QUINNS Training Generator Pro")
    
    # Render step indicator
    render_step_indicator(st.session_state.step)

    
    # STEP 1: Upload Documents
    if st.session_state.step == "upload":
       
        st.markdown("## 📤 Upload Training Documents")
        st.markdown("Upload one or more documents to generate comprehensive training materials.")
        
        uploaded_files = st.file_uploader(
            "Choose your documents",
            type=["pdf", "docx", "pptx", "txt", "md"],
            accept_multiple_files=True,
            help="Upload multiple documents - they will be combined for analysis"
        )
        
        if uploaded_files:
            st.markdown("### 📋 Uploaded Files")
            
            for idx, uploaded_file in enumerate(uploaded_files):
                filesize = f"{uploaded_file.size / 1024:.1f} KB" if uploaded_file.size < 1024*1024 else f"{uploaded_file.size / (1024*1024):.1f} MB"
                render_file_card(uploaded_file.name, filesize, idx)
            
            col1, col2 = st.columns([1, 4])
            with col1:
                if st.button("🚀 Process Documents", use_container_width=True):
                    if not st.session_state.api_key:
                        st.error("⚠️ API key not configured. Please set OPENAI_API_KEY environment variable or update the code with your API key.")
                    else:
                        st.session_state.file_paths = []
                        for uploaded_file in uploaded_files:
                            temp_path = TEMP_DIR / uploaded_file.name
                            with open(temp_path, "wb") as f:
                                f.write(uploaded_file.getvalue())
                            st.session_state.file_paths.append(temp_path)
                        
                        with st.spinner("🔄 Processing your documents..."):
                            try:
                                combined_content = ""
                                for file_path in st.session_state.file_paths:
                                    content = asyncio.run(st.session_state.generator.process_document(file_path))
                                    combined_content += f"\n\n=== Content from {file_path.name} ===\n\n{content}"
                                
                                st.session_state.combined_content = combined_content
                                
                                # Auto-detect training type if set to Auto-Detect
                                if st.session_state.training_type == "Auto-Detect":
                                    with st.spinner("🔍 Detecting training type..."):
                                        detected_type = st.session_state.generator.content_analyzer.detect_training_type(combined_content)
                                        st.session_state.training_type = detected_type
                                        st.success(f"✅ Detected training type: **{detected_type}**")
                                
                                if combined_content:
                                    st.success(f"✅ Successfully processed {len(uploaded_files)} document(s)")
                                    time.sleep(1)
                                    st.session_state.step = "topics"
                                    st.rerun()
                                else:
                                    st.error("❌ Could not extract content from documents")
                            except Exception as e:
                                st.error(f"❌ Error: {str(e)}")
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    # STEP 2: Review Topics (with quality analysis)
    elif st.session_state.step == "topics":
        
        st.markdown("## 🎯 Content Analysis & Topic Review")
        
        # Show detected/selected training type
        st.markdown(f"""
        <div class="info-box info">
            <strong>📋 Training Type:</strong> {st.session_state.training_type}<br>
            <strong>⏱️ Target Duration:</strong> {st.session_state.session_duration}<br>
            <strong>📊 Slide Configuration:</strong> ~{st.session_state.max_slides_per_module} slides per module
        </div>
        """, unsafe_allow_html=True)
        
        if not st.session_state.topics:
            with st.spinner("🔍 Analyzing content and extracting topics..."):
                try:
                    topics, outline = asyncio.run(
                        st.session_state.generator.analyze_content(
                            st.session_state.combined_content,
                            st.session_state.training_duration
                        )
                    )
                    st.session_state.topics = topics
                    st.session_state.outline = outline
                    
                    # Quality analysis
                    quality = asyncio.run(
                        st.session_state.generator.analyze_quality(
                            st.session_state.combined_content,
                            outline
                        )
                    )
                    st.session_state.quality_analysis = quality
                    
                except Exception as e:
                    st.error(f"❌ Error: {str(e)}")
                    if st.button("← Go Back"):
                        st.session_state.step = "upload"
                        st.rerun()
        
        if st.session_state.topics and st.session_state.quality_analysis:
            # Display quality metrics
            st.markdown("### Quality Analysis")
            
            qa = st.session_state.quality_analysis
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                st.markdown(f"""
                <div class="metric-card" style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);">
                    <div class="metric-label">Overall Score</div>
                    <div class="metric-value">{qa.get('overall_score', 0)}/100</div>
                </div>
                """, unsafe_allow_html=True)
            
            with col2:
                st.markdown(f"""
                <div class="metric-card" style="background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);">
                    <div class="metric-label">Readability</div>
                    <div class="metric-value">{qa.get('readability_score', 0)}/100</div>
                </div>
                """, unsafe_allow_html=True)
            
            with col3:
                st.markdown(f"""
                <div class="metric-card" style="background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);">
                    <div class="metric-label">Completeness</div>
                    <div class="metric-value">{qa.get('completeness_score', 0)}/100</div>
                </div>
                """, unsafe_allow_html=True)
            
            with col4:
                st.markdown(f"""
                <div class="metric-card" style="background: linear-gradient(135deg, #43e97b 0%, #38f9d7 100%);">
                    <div class="metric-label">Engagement</div>
                    <div class="metric-value">{qa.get('engagement_score', 0)}/100</div>
                </div>
                """, unsafe_allow_html=True)
            
            # Strengths and improvements
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("#### Strengths")
                for strength in qa.get('strengths', []):
                    st.markdown(f"- {strength}")
            
            with col2:
                st.markdown("#### Suggested Improvements")
                for improvement in qa.get('improvements', []):
                    st.markdown(f"- {improvement}")
            
            st.markdown("---")
            
            # Topic selection
            st.markdown("### Select Topics to Include")
            
            cols = st.columns(2)
            selected_topics = []
            
            for i, topic in enumerate(st.session_state.topics):
                with cols[i % 2]:
                    badge_class = f"badge-{topic.get('importance', 'medium')}"
                    
                    st.markdown(f"""
                    <div class="topic-card">
                        <strong>{topic.get('title', 'Topic')}</strong>
                        <span class="topic-badge {badge_class}">{topic.get('importance', 'medium').upper()}</span>
                        <div style="font-size: 0.875rem; color: #6B7280; margin-top: 0.5rem;">
                            {topic.get('description', '')}
                        </div>
                        <div style="font-size: 0.75rem; color: #9CA3AF; margin-top: 0.25rem;">
                            Est. {topic.get('estimated_duration_minutes', 30)} minutes
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    selected = st.checkbox(
                        f"Include this topic",
                        value=topic.get('importance') in ['high', 'medium'],
                        key=f"topic_{i}",
                        label_visibility="collapsed"
                    )
                    
                    if selected:
                        selected_topics.append(topic)
            
            st.markdown("---")
            
            col1, col2, col3 = st.columns([1, 1, 3])
            
            with col1:
                if st.button("← Back", use_container_width=True):
                    st.session_state.step = "upload"
                    st.rerun()
            
            with col2:
                if st.button("Next: Edit Outline →", use_container_width=True):
                    if selected_topics:
                        st.session_state.step = "edit"
                        st.rerun()
                    else:
                        st.warning("⚠️ Please select at least one topic")
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    # STEP 3: Edit Outline
    elif st.session_state.step == "edit":
        
        
        st.markdown("## ✏️ Edit Training Outline")
        
        # Estimate duration
        estimated_slides = sum(m.get('estimated_slides', 8) for m in st.session_state.outline.get('modules', []))
        duration_info = calculate_training_duration(estimated_slides, st.session_state.pace, st.session_state.include_assessments)
        
        st.markdown(f"""
        <div class="info-box info">
            <strong>Estimated Training Duration:</strong> {format_duration_display(duration_info)}<br>
            <strong>Estimated Slides:</strong> {estimated_slides}<br>
            <strong>Training Type:</strong> {st.session_state.training_type}
        </div>
        """, unsafe_allow_html=True)
        
        # Training Overview
        st.markdown("### Training Overview")
        
        col1, col2 = st.columns(2)
        
        with col1:
            outline_title = st.text_input(
                "Training Title",
                value=st.session_state.outline.get('title', 'Professional Training Program')
            )
            st.session_state.outline['title'] = outline_title
            
            target_audience = st.text_input(
                "Target Audience",
                value=st.session_state.outline.get('target_audience', 'Professionals')
            )
            st.session_state.outline['target_audience'] = target_audience
        
        with col2:
            description = st.text_area(
                "Description",
                value=st.session_state.outline.get('description', 'Comprehensive training program'),
                height=100
            )
            st.session_state.outline['description'] = description
        
        prerequisites = st.text_input(
            "Prerequisites",
            value=st.session_state.outline.get('prerequisites', 'Basic knowledge')
        )
        st.session_state.outline['prerequisites'] = prerequisites
        
        # Learning Objectives
        st.markdown("### Learning Objectives")
        
        objectives = st.session_state.outline.get('objectives', [])
        
        for idx, obj in enumerate(objectives):
            col1, col2 = st.columns([5, 1])
            with col1:
                new_obj = st.text_input(
                    f"Objective {idx + 1}",
                    value=obj,
                    key=f"obj_{idx}",
                    label_visibility="collapsed"
                )
                objectives[idx] = new_obj
            with col2:
                if st.button("Delete", key=f"del_obj_{idx}"):
                    objectives.pop(idx)
                    st.rerun()
        
        if st.button("Add Objective"):
            objectives.append("New learning objective")
            st.session_state.outline['objectives'] = objectives
            st.rerun()
        
        st.session_state.outline['objectives'] = objectives
        
        # Modules
        st.markdown("### Training Modules")
        
        modules = st.session_state.outline.get('modules', [])
        
        for mod_idx, module in enumerate(modules):
            with st.expander(f"Module {mod_idx + 1}: {module.get('title', 'Module')}", expanded=False):
                col1, col2 = st.columns(2)
                
                with col1:
                    module['title'] = st.text_input(
                        "Module Title",
                        value=module.get('title', ''),
                        key=f"mod_title_{mod_idx}"
                    )
                
                with col2:
                    module['duration'] = st.text_input(
                        "Duration",
                        value=module.get('duration', '45 minutes'),
                        key=f"mod_duration_{mod_idx}"
                    )
                
                # Module objectives
                st.markdown("**Module Objectives:**")
                mod_objectives = module.get('objectives', [])
                
                for obj_idx, obj in enumerate(mod_objectives):
                    col1, col2 = st.columns([5, 1])
                    with col1:
                        mod_objectives[obj_idx] = st.text_input(
                            f"Objective",
                            value=obj,
                            key=f"mod_{mod_idx}_obj_{obj_idx}",
                            label_visibility="collapsed"
                        )
                    with col2:
                        if st.button("Delete", key=f"del_mod_{mod_idx}_obj_{obj_idx}"):
                            mod_objectives.pop(obj_idx)
                            st.rerun()
                
                if st.button("Add Module Objective", key=f"add_mod_obj_{mod_idx}"):
                    mod_objectives.append("New module objective")
                    st.rerun()
                
                module['objectives'] = mod_objectives
                
                # Key Points
                st.markdown("**Key Points:**")
                key_points = module.get('key_points', [])
                
                for kp_idx, kp in enumerate(key_points):
                    col1, col2 = st.columns([5, 1])
                    with col1:
                        key_points[kp_idx] = st.text_input(
                            f"Key Point",
                            value=kp,
                            key=f"mod_{mod_idx}_kp_{kp_idx}",
                            label_visibility="collapsed"
                        )
                    with col2:
                        if st.button("Delete", key=f"del_mod_{mod_idx}_kp_{kp_idx}"):
                            key_points.pop(kp_idx)
                            st.rerun()
                
                if st.button("Add Key Point", key=f"add_mod_kp_{mod_idx}"):
                    key_points.append("New key point")
                    st.rerun()
                
                module['key_points'] = key_points
                
                # Delete module button
                if st.button(f"Delete Module {mod_idx + 1}", key=f"del_mod_{mod_idx}"):
                    modules.pop(mod_idx)
                    st.rerun()
        
        if st.button("Add New Module"):
            modules.append({
                "title": "New Module",
                "duration": "45 minutes",
                "objectives": ["Module objective"],
                "topics": ["Module topic"],
                "key_points": ["Key point"],
                "activities": ["Group discussion"],
                "estimated_slides": 8
            })
            st.session_state.outline['modules'] = modules
            st.rerun()
        
        st.session_state.outline['modules'] = modules
        
        st.markdown("---")
        
        col1, col2, col3 = st.columns([1, 1, 3])
        
        with col1:
            if st.button("Back to Topics", use_container_width=True):
                st.session_state.step = "topics"
                st.rerun()
        
        with col2:
            if st.button("Generate Materials", type="primary", use_container_width=True):
                st.session_state.step = "generate"
                st.rerun()
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    # STEP 4: Generate Materials
    elif st.session_state.step == "generate":
        
        st.markdown("## ⚙️ Generating Training Materials")
        
        output_dir = OUTPUT_DIR
        output_dir.mkdir(exist_ok=True)
        
        if 'generation_complete' not in st.session_state:
            st.session_state.generation_complete = False
            
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            try:
                status_text.info("Initializing generation process...")
                
                temp_content_file = TEMP_DIR / "combined_content.txt"
                with open(temp_content_file, "w", encoding="utf-8") as f:
                    f.write(st.session_state.combined_content)
                
                with st.spinner("Generating your training materials..."):
                    result = asyncio.run(
                        st.session_state.generator.generate_training_materials(
                            temp_content_file,
                            output_dir,
                            st.session_state.max_slides_per_module,
                            st.session_state.include_assessments,
                            st.session_state.training_duration
                        )
                    )
                    
                    st.session_state.output_files = result
                    st.session_state.generation_complete = True
                    
                    progress_bar.progress(100)
                    status_text.success("Training materials generated successfully!")
                    
                    time.sleep(1)
                    st.rerun()
                    
            except Exception as e:
                status_text.error(f"Error: {str(e)}")
                if st.button("Try Again"):
                    st.session_state.step = "edit"
                    st.rerun()
        
        if st.session_state.generation_complete:
            st.markdown('<div class="info-box success">', unsafe_allow_html=True)
            st.markdown("### Success! Your training materials are ready.")
            st.markdown('</div>', unsafe_allow_html=True)
            
            # Display generation summary
            duration_info = st.session_state.output_files.get('duration_info', {})
            total_slides = st.session_state.output_files.get('total_slides', 0)
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.markdown(f"""
                <div class="metric-card" style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);">
                    <div class="metric-label">Total Slides</div>
                    <div class="metric-value">{total_slides}</div>
                </div>
                """, unsafe_allow_html=True)
            
            with col2:
                st.markdown(f"""
                <div class="metric-card" style="background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);">
                    <div class="metric-label">Training Duration</div>
                    <div class="metric-value">{format_duration_display(duration_info)}</div>
                </div>
                """, unsafe_allow_html=True)
            
            with col3:
                st.markdown(f"""
                <div class="metric-card" style="background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);">
                    <div class="metric-label">Modules</div>
                    <div class="metric-value">{len(st.session_state.outline.get('modules', []))}</div>
                </div>
                """, unsafe_allow_html=True)
            
            st.markdown("### Download Your Materials")
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                if 'pptx_path' in st.session_state.output_files:
                    pptx_path = st.session_state.output_files['pptx_path']
                    if os.path.exists(pptx_path):
                        with open(pptx_path, "rb") as f:
                            st.download_button(
                                "Download PowerPoint",
                                f.read(),
                                file_name=Path(pptx_path).name,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True
                            )
            
            with col2:
                if 'guide_path' in st.session_state.output_files:
                    guide_path = st.session_state.output_files['guide_path']
                    if os.path.exists(guide_path):
                        with open(guide_path, "rb") as f:
                            st.download_button(
                                "Download Trainer Guide",
                                f.read(),
                                file_name=Path(guide_path).name,
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                use_container_width=True
                            )
            
            with col3:
                if st.session_state.include_assessments and st.session_state.output_files.get('assessment_path'):
                    assessment_path = st.session_state.output_files['assessment_path']
                    if os.path.exists(assessment_path):
                        with open(assessment_path, "rb") as f:
                            st.download_button(
                                "Download Assessment",
                                f.read(),
                                file_name=Path(assessment_path).name,
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                use_container_width=True
                            )
            
            st.markdown("---")
            
            if st.button("Start New Project", use_container_width=False):
                # Reset all session state
                keys_to_reset = [
                    'step', 'topics', 'outline', 'content', 'combined_content',
                    'output_files', 'file_paths', 'generation_complete',
                    'max_slides_per_module', 'quality_analysis', 'template',
                    'training_duration', 'include_assessments', 'pace'
                ]
                for key in keys_to_reset:
                    if key in st.session_state:
                        del st.session_state[key]
                
                st.session_state.step = "upload"
                st.session_state.max_slides_per_module = 12
                st.session_state.template = "corporate"
                st.session_state.training_duration = "1 day"
                st.session_state.include_assessments = True
                st.session_state.pace = "medium"
                st.rerun()
        
        st.markdown('</div>', unsafe_allow_html=True)

if __name__ == "__main__":
    main()